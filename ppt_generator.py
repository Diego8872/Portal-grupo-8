"""
ppt_generator.py — INTERLOG KPI Dashboard
Genera el PowerPoint con python-pptx replicando el diseño corporativo original.
Usa el KPI_template.pptx para extraer las imágenes de fondo.
"""

import io, os
from collections import Counter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FASA = 'FINNING ARGENTINA SOCIEDAD ANO'
FSM  = 'FINNING SOLUCIONES MINERAS SA'

# Colores exactos del template
TEAL    = RGBColor(0x3F, 0xA8, 0x9C)
DARK    = RGBColor(0x1E, 0x2D, 0x35)
GRAY    = RGBColor(0x88, 0x88, 0x88)
LGRAY   = RGBColor(0xCC, 0xCC, 0xCC)
BG_CARD = RGBColor(0xF0, 0xF4, 0xF6)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
VERDE   = RGBColor(0x27, 0xAE, 0x60)
NARANJA = RGBColor(0xD4, 0x69, 0x0A)
ROJO    = RGBColor(0xC0, 0x39, 0x2B)

# Colores tabla desvíos
ROW_DARK  = RGBColor(0xE8, 0xED, 0xF0)
ROW_LIGHT = RGBColor(0xF7, 0xF9, 0xFA)
HDR_TEAL  = RGBColor(0x2C, 0x7A, 0x72)

W_IN = 13.33
H_IN = 7.5

TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'KPI_template.pptx')


def calcular_kpi(items, con_parametros=True):
    total = len(items)
    if total == 0:
        return 0.0, 0, 0
    out = sum(1 for i in items
              if i['desvio'] and (not con_parametros or
                 str(i.get('parametro', '')).upper() == 'INTERLOG'))
    return round((total - out) / total * 100, 1), total - out, out


def kpi_color(pct):
    if pct >= 95: return TEAL
    if pct >= 80: return NARANJA
    return ROJO


def fmt_pct(pct):
    return f"{int(pct)}%" if pct == int(pct) else f"{pct:.1f}%"


def add_bg(slide, img_bytes):
    from io import BytesIO
    slide.shapes.add_picture(BytesIO(img_bytes), 0, 0,
                              Inches(W_IN), Inches(H_IN))


def rect(slide, x, y, w, h, fill, line=None, lw=0):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw) if lw else Pt(0.5)
    else:
        sh.line.fill.background()
    return sh


def txt(slide, text, x, y, w, h, size, bold=False, color=None, align=PP_ALIGN.LEFT,
        valign='top', italic=False):
    if color is None:
        color = DARK
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    from pptx.oxml.ns import qn as _qn
    txBody = tf._txBody
    bodyPr = txBody.find(_qn('a:bodyPr'))
    if bodyPr is None:
        import lxml.etree as _et
        bodyPr = _et.SubElement(txBody, _qn('a:bodyPr'))
    bodyPr.set('anchor', {'middle': 'ctr', 'bottom': 'b'}.get(valign, 't'))
    return tb


def slide_header(slide, title):
    txt(slide, title, 0.45, 0.30, 12.0, 0.55, 22, bold=True, color=DARK)
    rect(slide, 0.45, 0.88, 3.0, 0.04, TEAL)


def add_kpi_card(slide, x, total, inn, out, pct, title, subtitle, limite):
    color = kpi_color(pct)
    w_card = 2.80

    rect(slide, x, 1.05, w_card, 5.70, BG_CARD)
    rect(slide, x, 1.05, 0.06, 5.70, TEAL)

    txt(slide, title,    x+0.15, 1.10, 2.5, 0.35, 14, bold=True,  color=DARK)
    txt(slide, subtitle, x+0.15, 1.44, 2.5, 0.25,  8, color=GRAY)
    rect(slide, x+0.20, 1.95, 2.40, 0.02, LGRAY)
    txt(slide, 'TARGET', x+0.15, 2.05, 2.5, 0.25,  8, color=GRAY)
    txt(slide, '95%',    x+0.15, 2.28, 2.5, 0.30, 14, bold=True, color=DARK)
    txt(slide, 'LÍMITE', x+0.15, 2.68, 2.5, 0.25,  8, color=GRAY)
    txt(slide, limite,   x+0.15, 2.90, 2.5, 0.30, 14, bold=True, color=DARK)
    rect(slide, x+0.20, 3.60, 2.40, 0.02, LGRAY)

    for i, (v, lbl, col) in enumerate([(str(total), 'TOTAL', DARK),
                                         (str(inn),   'IN',    TEAL),
                                         (str(out),   'OUT',   ROJO if out > 0 else TEAL)]):
        sx = x + 0.15 + i * 0.80
        txt(slide, v,   sx, 3.70, 0.75, 0.42, 20, bold=True, color=col)
        txt(slide, lbl, sx, 4.12, 0.75, 0.25,  9, color=GRAY)

    kx = x + w_card + 0.35
    kw = 2.90
    txt(slide, fmt_pct(pct), kx, 1.80, kw, 1.60,
        52, bold=True, color=color, align=PP_ALIGN.CENTER, valign='middle')
    txt(slide, 'KPI IN', kx, 3.40, kw, 0.28,
        16, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    rect(slide, kx+0.1, 4.05, 2.90, 0.28, BG_CARD)
    fw = max(0.05, 2.90 * pct / 100)
    rect(slide, kx+0.1, 4.05, fw, 0.28, color)
    txt(slide, 'TARGET 95%', kx, 4.38, kw, 0.28,
        16, bold=True, color=GRAY, align=PP_ALIGN.CENTER)


def add_liberacion_slide(prs, bg_bytes, lib_items, nombre, blank_layout):
    s = prs.slides.add_slide(blank_layout)
    add_bg(s, bg_bytes)
    slide_header(s, f"{nombre} · LIBERACIÓN")

    sub = [i for i in lib_items if i['nombre'] == nombre]
    pct_tot, in_tot, out_tot = calcular_kpi(sub, True)

    txt(s, 'RESUMEN', 0.45, 0.98, 3.0, 0.24, 9, bold=True, color=TEAL)

    boxes = [(str(len(sub)), 'TOTAL OPS', DARK),
             (str(in_tot),  'IN',        TEAL),
             (str(out_tot), 'OUT',        ROJO if out_tot > 0 else TEAL),
             (fmt_pct(pct_tot), 'KPI',    kpi_color(pct_tot))]
    for i, (v, lbl, col) in enumerate(boxes):
        bx = 0.45 + i * 1.65
        rect(s, bx, 1.22, 1.50, 0.88, WHITE)
        txt(s, v,   bx, 1.26, 1.50, 0.52, 26, bold=True, color=col,
            align=PP_ALIGN.CENTER, valign='middle')
        txt(s, lbl, bx, 1.78, 1.50, 0.25,  8, bold=True, color=GRAY,
            align=PP_ALIGN.CENTER)

    txt(s, 'DETALLE POR VÍA', 0.45, 2.25, 5.0, 0.28, 9, bold=True, color=TEAL)

    via_labels    = {'AVION': 'Aéreo', 'CAMION': 'Camión', 'MARITIMO': 'Marítimo'}
    canal_colores = {'VERDE': VERDE, 'NARANJA': NARANJA, 'ROJO': ROJO}

    row_y = 2.58
    row_h = 1.30
    row_gap = 0.10

    for via in ['AVION', 'CAMION', 'MARITIMO']:
        via_items = [i for i in sub if i['via'] == via]
        if not via_items:
            continue

        pct_v, in_v, out_v = calcular_kpi(via_items, True)
        total_v = len(via_items)
        color_v = kpi_color(pct_v)

        rect(s, 0.45, row_y, 12.40, row_h, BG_CARD)
        txt(s, via_labels[via], 0.55, row_y+0.06, 1.80, 0.28, 10, bold=True, color=DARK)
        txt(s, f"Total: {total_v} · {in_v} IN · {out_v} OUT",
            0.55, row_y+0.34, 2.00, 0.20, 7.5, color=GRAY)
        txt(s, fmt_pct(pct_v), 2.55, row_y+0.06, 1.00, 0.50,
            18, bold=True, color=color_v, align=PP_ALIGN.CENTER, valign='middle')

        bar_x0 = 3.65
        bw_each = 2.90
        bgap    = 0.22

        for ci, canal in enumerate(['VERDE', 'NARANJA', 'ROJO']):
            canal_items = [i for i in via_items if i['canal'] == canal]
            cnt = len(canal_items)
            pct_c = round(cnt / total_v * 100) if total_v > 0 else 0
            cc  = canal_colores[canal]
            cx  = bar_x0 + ci * (bw_each + bgap)

            txt(s, canal, cx, row_y+0.06, bw_each, 0.20, 7.5, bold=True, color=cc)
            rect(s, cx, row_y+0.30, bw_each, 0.28, LGRAY)
            if cnt > 0:
                fw = max(0.05, bw_each * pct_c / 100)
                rect(s, cx, row_y+0.30, fw, 0.28, cc)
            val_str = f"{cnt}/{total_v}   {pct_c}%" if cnt > 0 else "–"
            txt(s, val_str, cx, row_y+0.62, bw_each, 0.20, 8, color=GRAY)

        row_y += row_h + row_gap

    txt(s, 'Target: 95%', 7.20, 1.22, 5.00, 0.88,
        30, bold=True, color=DARK, align=PP_ALIGN.CENTER, valign='middle')


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE DE DESVÍOS  (funciones unificadas y auto-ajustables)
# ═══════════════════════════════════════════════════════════════════════════

# Dimensiones del slide (pulgadas)
_SLIDE_W  = W_IN        # 13.33
_SLIDE_H  = H_IN        # 7.50
_MARGIN   = 0.38        # margen lateral
_FOOTER_H = 0.60        # alto reservado para footer + nota al pie
_HDR_TOP  = 1.00        # y donde empieza el contenido tras el header


def _build_filas(items, es_cm=False):
    """Construye filas para la tabla de desvíos de un grupo de items."""
    total = len(items)
    pct_kpi, in_count, out_count = calcular_kpi(items, True)
    sin_desvio = sum(1 for i in items if not i['desvio'])

    param_counter = Counter()
    for i in items:
        if i['desvio']:
            param = str(i.get('parametro', '')).strip().upper()
            param_counter[param if param else 'PENDIENTE'] += 1

    filas = []
    pct_in = round(in_count / total * 100, 2) if total else 0
    filas.append(('IN', in_count, f"{pct_in:.2f}%", TEAL, True))
    interlog_cnt = param_counter.pop('INTERLOG', 0)
    if interlog_cnt:
        pct_p = round(interlog_cnt / total * 100, 2) if total else 0
        filas.append(('INTERLOG', interlog_cnt, f"{pct_p:.2f}%", ROJO, False))
    for param, cnt in sorted(param_counter.items()):
        pct_p = round(cnt / total * 100, 2) if total else 0
        filas.append((param, cnt, f"{pct_p:.2f}%", NARANJA, False))
    pct_sd = round(sin_desvio / total * 100, 2) if total else 0
    filas.append(('S/DESVÍO', sin_desvio, f"{pct_sd:.2f}%", GRAY, False))
    filas.append(('Total general', total, '100,00%', DARK, True))
    unit = 'exp' if es_cm else 'ops'
    return filas, pct_kpi, in_count, out_count, total, unit


def _draw_desvios_tabla(s, items_data, titulo_col,
                        x, tbl_w, kpi_y, hdr_y, row_h, font_row,
                        es_cm=False):
    """
    Dibuja una columna completa de desvíos (KPI + tabla + barra) en la posición dada.
    items_data: ya pre-calculado como (filas, pct_kpi, in_count, out_count, total, unit)
    """
    filas, pct_kpi, in_count, out_count, total, unit = items_data
    kpi_col = kpi_color(pct_kpi)

    # ── Proporciones internas de columnas
    col_rango = tbl_w * 0.58
    col_op    = tbl_w * 0.17
    col_kpi_c = tbl_w * 0.25
    font_hdr  = max(8.0, font_row)

    # ── Etiqueta sociedad / título de columna
    if titulo_col:
        txt(s, titulo_col, x, kpi_y - 0.28, tbl_w, 0.26,
            10, bold=True, color=TEAL)

    # ── KPI grande
    txt(s, fmt_pct(pct_kpi), x, kpi_y, tbl_w, 0.72,
        32, bold=True, color=kpi_col, align=PP_ALIGN.CENTER, valign='middle')
    txt(s, f"{in_count} IN  ·  {out_count} OUT  ·  {total} {unit}",
        x, kpi_y + 0.72, tbl_w, 0.22,
        7.5, color=GRAY, align=PP_ALIGN.CENTER)

    # ── Header tabla
    hdr_h = 0.36
    rect(s, x,                       hdr_y, col_rango, hdr_h, HDR_TEAL)
    rect(s, x + col_rango,           hdr_y, col_op,    hdr_h, HDR_TEAL)
    rect(s, x + col_rango + col_op,  hdr_y, col_kpi_c, hdr_h, HDR_TEAL)
    txt(s, 'Rango', x + 0.10, hdr_y, col_rango - 0.10, hdr_h,
        font_hdr, bold=True, color=WHITE, valign='middle')
    txt(s, 'OP', x + col_rango + 0.04, hdr_y, col_op - 0.04, hdr_h,
        font_hdr, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign='middle')
    txt(s, 'KPI', x + col_rango + col_op + 0.04, hdr_y, col_kpi_c - 0.04, hdr_h,
        font_hdr, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign='middle')

    # ── Filas de datos
    for ri, (lbl, ops, kpi_str, col_val, is_bold) in enumerate(filas):
        ry = hdr_y + hdr_h + ri * row_h
        bg = ROW_DARK if ri % 2 == 0 else ROW_LIGHT
        if lbl == 'Total general':
            bg = RGBColor(0xD5, 0xE5, 0xE3)

        rect(s, x,                      ry, col_rango, row_h, bg)
        rect(s, x + col_rango,          ry, col_op,    row_h, bg)
        rect(s, x + col_rango + col_op, ry, col_kpi_c, row_h, bg)

        if lbl in ['IN', 'Total general']:
            rect(s, x, ry, 0.05, row_h, TEAL if lbl == 'IN' else HDR_TEAL)

        c = DARK if lbl in ['S/DESVÍO', 'Total general'] else col_val
        txt(s, lbl,
            x + 0.10, ry, col_rango - 0.10, row_h,
            font_row, bold=is_bold, color=c, valign='middle')
        txt(s, str(ops),
            x + col_rango + 0.04, ry, col_op - 0.04, row_h,
            font_row, bold=is_bold, color=c, align=PP_ALIGN.CENTER, valign='middle')
        txt(s, kpi_str,
            x + col_rango + col_op + 0.04, ry, col_kpi_c - 0.04, row_h,
            font_row, bold=is_bold, color=c, align=PP_ALIGN.CENTER, valign='middle')

    # ── Barra de progreso
    bar_y = hdr_y + hdr_h + len(filas) * row_h + 0.10
    bar_h = 0.16
    rect(s, x, bar_y, tbl_w, bar_h, LGRAY)
    bw = max(0.05, tbl_w * pct_kpi / 100)
    rect(s, x, bar_y, bw, bar_h, kpi_col)
    txt(s, 'TARGET: 95%', x, bar_y + bar_h + 0.05, tbl_w, 0.20,
        7.5, color=GRAY, align=PP_ALIGN.CENTER)


def _calc_layout_doble(max_filas):
    """
    Calcula row_h y font_row para slides de dos columnas,
    garantizando que todo quede dentro del slide.
    """
    GAP      = 0.25
    SLIDE_W  = _SLIDE_W - 2 * _MARGIN          # espacio total disponible
    tbl_w    = (SLIDE_W - GAP) / 2              # ancho de cada columna

    # Espacio vertical disponible para las filas
    kpi_blk  = 1.16   # etiqueta + kpi grande + contador
    hdr_h    = 0.36
    bar_blk  = 0.16 + 0.05 + 0.20 + 0.10       # barra + espacios
    nota_blk = _FOOTER_H
    hdr_y    = _HDR_TOP + kpi_blk + 0.14        # donde empieza el header

    available = (_SLIDE_H - hdr_y - hdr_h - bar_blk - nota_blk)
    row_h     = min(0.44, available / max_filas) if max_filas else 0.44
    font_row  = max(7.0, min(9.5, row_h * 20))

    x_left  = _MARGIN
    x_right = _MARGIN + tbl_w + GAP
    kpi_y   = _HDR_TOP + 0.30

    return tbl_w, row_h, font_row, hdr_y, kpi_y, x_left, x_right


def _calc_layout_simple(max_filas):
    """
    Calcula layout para slide de una sola columna + panel KPI derecho.
    """
    KPI_PANEL_W = 3.30
    tbl_x   = _MARGIN
    tbl_w   = _SLIDE_W - _MARGIN - KPI_PANEL_W - 0.50  # ≈ 9.10
    hdr_h   = 0.42
    hdr_y   = _HDR_TOP + 0.05
    bar_blk = 0.42
    nota_blk= _FOOTER_H

    available = (_SLIDE_H - hdr_y - hdr_h - bar_blk - nota_blk)
    row_h     = min(0.46, available / max_filas) if max_filas else 0.46
    font_row  = max(7.5, min(10.5, row_h * 21))

    kpi_x = _SLIDE_W - _MARGIN - KPI_PANEL_W
    return tbl_x, tbl_w, hdr_y, hdr_h, row_h, font_row, kpi_x, KPI_PANEL_W


# ── API pública ─────────────────────────────────────────────────────────────

def add_desvios_doble_slide(prs, bg_bytes, fasa_items, fsm_items,
                            titulo_slide, blank_layout, es_cm=False):
    """
    Slide con DOS columnas: FASA (izquierda) y FSM (derecha).
    Se usa para Oficialización y Canal Verde.
    Layout completamente auto-ajustable al número de parámetros de desvío.
    """
    s = prs.slides.add_slide(blank_layout)
    add_bg(s, bg_bytes)
    slide_header(s, titulo_slide)

    data_fasa = _build_filas(fasa_items, es_cm)
    data_fsm  = _build_filas(fsm_items,  es_cm) if fsm_items else None

    max_filas = max(
        len(data_fasa[0]),
        len(data_fsm[0]) if data_fsm else 0
    )
    tbl_w, row_h, font_row, hdr_y, kpi_y, x_left, x_right = _calc_layout_doble(max_filas)

    _draw_desvios_tabla(s, data_fasa, 'FASA',
                        x_left, tbl_w, kpi_y, hdr_y, row_h, font_row, es_cm)

    # Separador
    sep_x = x_left + tbl_w + (x_right - x_left - tbl_w) / 2 - 0.01
    rect(s, sep_x, _HDR_TOP, 0.02, _SLIDE_H - _HDR_TOP - _FOOTER_H + 0.20, LGRAY)

    if data_fsm:
        _draw_desvios_tabla(s, data_fsm, 'FSM',
                            x_right, tbl_w, kpi_y, hdr_y, row_h, font_row, es_cm)

    nota_y = _SLIDE_H - _FOOTER_H + 0.05
    txt(s, '* Parámetro ≠ INTERLOG → operación se considera IN  |  S/DESVÍO = dentro del rango',
        _MARGIN, nota_y, _SLIDE_W - 2 * _MARGIN, 0.28, 8, color=GRAY, italic=True)


def add_desvios_canal_verde_slide(prs, bg_bytes, lib_items, blank_layout):
    fasa = [i for i in lib_items if i['via'] == 'AVION' and i['canal'] == 'VERDE' and i['nombre'] == 'FASA']
    fsm  = [i for i in lib_items if i['via'] == 'AVION' and i['canal'] == 'VERDE' and i['nombre'] == 'FSM']
    add_desvios_doble_slide(prs, bg_bytes, fasa, fsm,
                            'DESVÍOS · CANAL VERDE VÍA AÉREA', blank_layout)


def add_desvios_slide(prs, bg_bytes, items, titulo, blank_layout, es_cm=False):
    """
    Slide de desvíos para un único grupo (Liberación FASA, Liberación FSM, CM).
    Tabla a la izquierda + panel KPI derecho. Auto-ajustable.
    """
    s = prs.slides.add_slide(blank_layout)
    add_bg(s, bg_bytes)
    slide_header(s, titulo)

    data = _build_filas(items, es_cm)
    filas, pct_kpi, in_count, out_count, total, unit = data
    max_filas = len(filas)

    tbl_x, tbl_w, hdr_y, hdr_h, row_h, font_row, kpi_x, KPI_PANEL_W = \
        _calc_layout_simple(max_filas)

    kpi_col = kpi_color(pct_kpi)

    # ── Panel KPI derecho ─────────────────────────────────────────────────
    panel_h = 2.60
    rect(s, kpi_x, 1.05, KPI_PANEL_W, panel_h, BG_CARD)
    rect(s, kpi_x, 1.05, 0.06, panel_h, kpi_col)
    txt(s, 'KPI IN', kpi_x + 0.12, 1.12, KPI_PANEL_W - 0.18, 0.28,
        10, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    txt(s, fmt_pct(pct_kpi), kpi_x + 0.12, 1.38, KPI_PANEL_W - 0.18, 1.20,
        50, bold=True, color=kpi_col, align=PP_ALIGN.CENTER, valign='middle')
    txt(s, 'TARGET: 95%', kpi_x + 0.12, 2.62, KPI_PANEL_W - 0.18, 0.25,
        9, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    rect(s, kpi_x + 0.18, 2.92, KPI_PANEL_W - 0.36, 0.20, LGRAY)
    bw_prog = max(0.05, (KPI_PANEL_W - 0.36) * pct_kpi / 100)
    rect(s, kpi_x + 0.18, 2.92, bw_prog, 0.20, kpi_col)
    txt(s, f"{in_count} IN  ·  {out_count} OUT",
        kpi_x + 0.12, 3.18, KPI_PANEL_W - 0.18, 0.28,
        9, color=GRAY, align=PP_ALIGN.CENTER)
    txt(s, f"Total: {total} {unit}",
        kpi_x + 0.12, 3.46, KPI_PANEL_W - 0.18, 0.25,
        9, color=GRAY, align=PP_ALIGN.CENTER)

    # ── Tabla de desvíos ──────────────────────────────────────────────────
    col_rango = tbl_w * 0.57
    col_op    = tbl_w * 0.18
    col_kpi_c = tbl_w * 0.25
    font_hdr  = max(8.5, font_row)

    # Header
    rect(s, tbl_x,                       hdr_y, col_rango, hdr_h, HDR_TEAL)
    rect(s, tbl_x + col_rango,           hdr_y, col_op,    hdr_h, HDR_TEAL)
    rect(s, tbl_x + col_rango + col_op,  hdr_y, col_kpi_c, hdr_h, HDR_TEAL)
    txt(s, 'Rango',
        tbl_x + 0.12, hdr_y, col_rango - 0.12, hdr_h,
        font_hdr, bold=True, color=WHITE, valign='middle')
    txt(s, 'OP',
        tbl_x + col_rango + 0.06, hdr_y, col_op - 0.06, hdr_h,
        font_hdr, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign='middle')
    txt(s, 'KPI',
        tbl_x + col_rango + col_op + 0.06, hdr_y, col_kpi_c - 0.06, hdr_h,
        font_hdr, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign='middle')

    # Filas
    for ri, (lbl, ops, kpi_str, col_val, is_bold) in enumerate(filas):
        ry = hdr_y + hdr_h + ri * row_h
        bg = ROW_DARK if ri % 2 == 0 else ROW_LIGHT
        if lbl == 'Total general':
            bg = RGBColor(0xD5, 0xE5, 0xE3)

        rect(s, tbl_x,                      ry, col_rango, row_h, bg)
        rect(s, tbl_x + col_rango,          ry, col_op,    row_h, bg)
        rect(s, tbl_x + col_rango + col_op, ry, col_kpi_c, row_h, bg)

        if lbl in ['IN', 'Total general']:
            rect(s, tbl_x, ry, 0.06, row_h, TEAL if lbl == 'IN' else HDR_TEAL)

        c = DARK if lbl in ['S/DESVÍO', 'Total general'] else col_val
        txt(s, lbl,
            tbl_x + 0.14, ry, col_rango - 0.14, row_h,
            font_row, bold=is_bold, color=c, valign='middle')
        txt(s, str(ops),
            tbl_x + col_rango + 0.06, ry, col_op - 0.06, row_h,
            font_row, bold=is_bold, color=c, align=PP_ALIGN.CENTER, valign='middle')
        txt(s, kpi_str,
            tbl_x + col_rango + col_op + 0.06, ry, col_kpi_c - 0.06, row_h,
            font_row, bold=is_bold, color=c, align=PP_ALIGN.CENTER, valign='middle')

    # Nota al pie
    nota_y = hdr_y + hdr_h + len(filas) * row_h + 0.16
    txt(s, '* Parámetro ≠ INTERLOG → operación se considera IN  |  S/DESVÍO = dentro del rango',
        tbl_x, nota_y, tbl_w + 0.50, 0.28, 8, color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════
def generar_ppt(lib_items, ofi_items, cm_pre_items, cm_apr_items, mes='MES'):
    """Genera el PowerPoint y retorna BytesIO."""

    def kpi_nom(items, nombre):
        sub = [i for i in items if i['nombre'] == nombre]
        pct, inn, out = calcular_kpi(sub, True)
        return pct, inn, out, len(sub)

    def kpi_cv(nombre):
        sub = [i for i in lib_items
               if i['nombre'] == nombre and i['via'] == 'AVION' and i['canal'] == 'VERDE']
        if not sub: return 0.0, 0, 0, 0
        pct, inn, out = calcular_kpi(sub, True)
        return pct, inn, out, len(sub)

    def dev_imp(items, nombre):
        sub = [i for i in items if i['nombre'] == nombre]
        out = sum(1 for i in sub if i['desvio'] and
                  str(i.get('parametro', '')).upper() == 'INTERLOG')
        return out, len(sub)

    lib_fasa_pct, lib_fasa_in, lib_fasa_out, lib_fasa_tot = kpi_nom(lib_items, 'FASA')
    lib_fsm_pct,  lib_fsm_in,  lib_fsm_out,  lib_fsm_tot  = kpi_nom(lib_items, 'FSM')
    ofi_fasa_pct, ofi_fasa_in, ofi_fasa_out, ofi_fasa_tot = kpi_nom(ofi_items, 'FASA')
    ofi_fsm_pct,  ofi_fsm_in,  ofi_fsm_out,  ofi_fsm_tot  = kpi_nom(ofi_items, 'FSM')
    cv_fasa_pct,  cv_fasa_in,  cv_fasa_out,  cv_fasa_tot  = kpi_cv('FASA')
    cv_fsm_pct,   cv_fsm_in,   cv_fsm_out,   cv_fsm_tot   = kpi_cv('FSM')
    cm_pct, cm_in, cm_out = calcular_kpi(cm_pre_items, True)
    cm_tot = len(cm_pre_items)

    rangos  = Counter(i['rango'] for i in cm_apr_items if i.get('rango'))
    tot_apr = sum(rangos.values())
    r0_7    = rangos.get('0 a 7', 0)
    r8_15   = rangos.get('8 a 15', 0)
    r15     = rangos.get('+15', 0)
    def pct_r(v): return round(v / tot_apr * 100) if tot_apr else 0

    candidates = [TEMPLATE_PATH,
                  '/mnt/user-data/uploads/KPI_INTERLOG_Diciembre_2025__21_.pptx']
    tfile = next((p for p in candidates if os.path.exists(p)), None)
    if not tfile:
        raise FileNotFoundError("No se encontró KPI_template.pptx")

    tprs = Presentation(tfile)
    bg_portada = bg_contenido = None
    for idx, sl in enumerate(tprs.slides):
        for sh in sl.shapes:
            if sh.shape_type == 13:
                blob = sh.image.blob
                if idx == 0:
                    bg_portada = blob
                else:
                    bg_contenido = blob
                break
        if bg_portada and bg_contenido:
            break

    prs = Presentation()
    prs.slide_width  = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    blank = prs.slide_layouts[6]

    # ── SLIDE 1: PORTADA ────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    add_bg(s1, bg_portada)
    txt(s1, 'KPI FASA / FSM', 9.00, 2.55, 4.00, 0.60,
        24, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(s1, mes, 9.00, 3.18, 4.00, 0.45,
        20, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(s1, 'Reporte de Indicadores de Desempeño', 9.00, 3.68, 4.00, 0.28,
        10, color=DARK, align=PP_ALIGN.CENTER)

    # ── SLIDE 2: RESUMEN EJECUTIVO ──────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    add_bg(s2, bg_contenido)
    slide_header(s2, 'RESUMEN EJECUTIVO')

    kpis = [
        (fmt_pct(lib_fasa_pct), 'LIBERACIÓN', 'FASA', f"{lib_fasa_tot} operaciones", kpi_color(lib_fasa_pct)),
        (fmt_pct(lib_fsm_pct),  'LIBERACIÓN', 'FSM',  f"{lib_fsm_tot} operaciones",  kpi_color(lib_fsm_pct)),
        (fmt_pct(ofi_fasa_pct), 'OFICIALIZACIÓN', 'FASA', f"{ofi_fasa_tot} operaciones", kpi_color(ofi_fasa_pct)),
        (fmt_pct(ofi_fsm_pct),  'OFICIALIZACIÓN', 'FSM',  f"{ofi_fsm_tot} operaciones",  kpi_color(ofi_fsm_pct)),
        (fmt_pct(cm_pct),       'CM', 'PRESENTADOS', f"{cm_tot} expedientes", kpi_color(cm_pct)),
    ]
    cxs = [0.45, 2.85, 5.25, 7.65, 10.05]
    cw  = 2.20
    for (pct_str, l1, l2, sub, col), cx in zip(kpis, cxs):
        rect(s2, cx, 1.05, cw, 1.85, WHITE)
        rect(s2, cx, 1.05, 0.05, 1.85, col)
        txt(s2, pct_str, cx, 1.10, cw, 0.78, 28, bold=True, color=col,
            align=PP_ALIGN.CENTER, valign='middle')
        txt(s2, l1, cx, 1.88, cw, 0.25,  9, bold=True, color=DARK,
            align=PP_ALIGN.CENTER)
        txt(s2, l2, cx, 2.12, cw, 0.22,  9, bold=True, color=DARK,
            align=PP_ALIGN.CENTER)
        txt(s2, sub, cx, 2.34, cw, 0.22, 7.5, color=GRAY, align=PP_ALIGN.CENTER)

    txt(s2, 'Target: 95%  ·  Parámetro ≠ INTERLOG → operación IN',
        0.45, 3.08, 12.0, 0.35, 16, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    txt(s2, 'DESVÍOS IMPUTABLES A INTERLOG', 0.45, 3.58, 8.0, 0.28,
        9, bold=True, color=TEAL)

    cm_imp = sum(1 for i in cm_pre_items if i['desvio'] and
                 str(i.get('parametro', '')).upper() == 'INTERLOG')
    devs = [('FASA · Liberación',     *dev_imp(lib_items, 'FASA')),
            ('FSM · Liberación',      *dev_imp(lib_items, 'FSM')),
            ('FASA · Oficialización', *dev_imp(ofi_items, 'FASA')),
            ('FSM · Oficialización',  *dev_imp(ofi_items, 'FSM')),
            ('CM Presentados',        cm_imp, len(cm_pre_items))]
    dxs = [0.45, 2.90, 5.35, 7.80, 10.25]
    dw  = 2.30
    for (lbl, dout, dtot), dx in zip(devs, dxs):
        rect(s2, dx, 3.85, dw, 1.10, WHITE)
        txt(s2, lbl, dx+0.10, 3.90, dw-0.15, 0.25, 8, color=GRAY)
        txt(s2, f"{dout} OUT", dx+0.10, 4.16, dw*0.60, 0.44,
            18, bold=True, color=ROJO if dout > 0 else TEAL)
        unit = 'exp' if lbl == 'CM Presentados' else 'ops'
        txt(s2, f"de {dtot} {unit}", dx+0.10, 4.62, dw-0.15, 0.24, 9, color=GRAY)

    # ── SLIDE 3: OFICIALIZACIÓN · FASA + FSM ────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    add_bg(s3, bg_contenido)
    slide_header(s3, 'OFICIALIZACIÓN · FASA + FSM')
    add_kpi_card(s3, 0.45, ofi_fasa_tot, ofi_fasa_in, ofi_fasa_out, ofi_fasa_pct,
                 'FASA', 'OFICIALIZACIÓN', '1 día hábil')
    add_kpi_card(s3, 6.80, ofi_fsm_tot,  ofi_fsm_in,  ofi_fsm_out,  ofi_fsm_pct,
                 'FSM',  'OFICIALIZACIÓN', '1 día hábil (2 días Marítimo)')

    # ── SLIDE 4: DESVÍOS OFICIALIZACIÓN (FASA + FSM separados) ──────────────
    ofi_fasa = [i for i in ofi_items if i['nombre'] == 'FASA']
    ofi_fsm  = [i for i in ofi_items if i['nombre'] == 'FSM']
    add_desvios_doble_slide(prs, bg_contenido, ofi_fasa, ofi_fsm,
                            'DESVÍOS · OFICIALIZACIÓN FASA + FSM', blank)

    # ── SLIDE 5: FASA · LIBERACIÓN ──────────────────────────────────────────
    add_liberacion_slide(prs, bg_contenido, lib_items, 'FASA', blank)

    # ── SLIDE 6: DESVÍOS LIBERADAS FASA ─────────────────────────────────────
    lib_fasa_items = [i for i in lib_items if i['nombre'] == 'FASA']
    add_desvios_slide(prs, bg_contenido, lib_fasa_items,
                      'DESVÍOS · LIBERACIÓN FASA', blank)

    # ── SLIDE 7: FSM · LIBERACIÓN ───────────────────────────────────────────
    add_liberacion_slide(prs, bg_contenido, lib_items, 'FSM', blank)

    # ── SLIDE 8: DESVÍOS LIBERADAS FSM ──────────────────────────────────────
    lib_fsm_items = [i for i in lib_items if i['nombre'] == 'FSM']
    add_desvios_slide(prs, bg_contenido, lib_fsm_items,
                      'DESVÍOS · LIBERACIÓN FSM', blank)

    # ── SLIDE 9: CANAL VERDE · VÍA AÉREA ────────────────────────────────────
    s9 = prs.slides.add_slide(blank)
    add_bg(s9, bg_contenido)
    slide_header(s9, 'CANAL VERDE · VÍA AÉREA')
    add_kpi_card(s9, 0.45, cv_fasa_tot, cv_fasa_in, cv_fasa_out, cv_fasa_pct,
                 'FASA', 'CANAL VERDE AÉREO', '1 día hábil')
    add_kpi_card(s9, 6.80, cv_fsm_tot,  cv_fsm_in,  cv_fsm_out,  cv_fsm_pct,
                 'FSM',  'CANAL VERDE AÉREO', '1 día hábil')

    # ── SLIDE 10: DESVÍOS CANAL VERDE VÍA AÉREA ─────────────────────────────
    add_desvios_canal_verde_slide(prs, bg_contenido, lib_items, blank)

    # ── SLIDE 11: DISTRIBUCIÓN DE CANALES ────────────────────────────────────
    s11 = prs.slides.add_slide(blank)
    add_bg(s11, bg_contenido)
    slide_header(s11, 'DISTRIBUCIÓN DE CANALES · LIBERACIONES')

    rows = []
    for nombre in ['FASA', 'FSM']:
        for via in ['AVION', 'CAMION', 'MARITIMO']:
            via_items = [i for i in lib_items if i['nombre'] == nombre and i['via'] == via]
            if not via_items: continue
            via_lbl = {'AVION': 'Aéreo', 'CAMION': 'Camión', 'MARITIMO': 'Marítimo'}[via]
            rows.append({
                'label':   f"{nombre} - {via_lbl}",
                'sub':     f"{len(via_items)} ops",
                'total':   len(via_items),
                'verde':   sum(1 for i in via_items if i['canal'] == 'VERDE'),
                'naranja': sum(1 for i in via_items if i['canal'] == 'NARANJA'),
                'rojo':    sum(1 for i in via_items if i['canal'] == 'ROJO'),
            })

    bar_x    = 2.60
    bar_w    = 9.50
    rh       = 0.62
    ry_start = 1.10
    rgap     = 0.18

    for ri, row in enumerate(rows):
        ry = ry_start + ri * (rh + rgap)
        total_r = row['total']
        txt(s11, row['label'], 0.45, ry, bar_x-0.55, rh*0.55, 10, bold=True, color=DARK)
        txt(s11, row['sub'],   0.45, ry+rh*0.52, bar_x-0.55, rh*0.40, 8, color=GRAY)
        if total_r == 0:
            rect(s11, bar_x, ry+0.08, bar_w, rh-0.16, BG_CARD)
            continue
        cx = bar_x
        for cnt, col in [(row['verde'], VERDE), (row['naranja'], NARANJA), (row['rojo'], ROJO)]:
            if cnt == 0: continue
            sw = bar_w * cnt / total_r
            rect(s11, cx, ry+0.08, sw, rh-0.16, col)
            pct_seg = round(cnt / total_r * 100)
            if sw > 0.5:
                txt(s11, f"{cnt}({pct_seg}%)", cx, ry+0.08, sw, rh-0.16,
                    9, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign='middle')
            cx += sw

    leg_y = 6.50
    for li, (lbl, col) in enumerate([('Canal Verde', VERDE), ('Canal Naranja', NARANJA), ('Canal Rojo', ROJO)]):
        lx = 2.20 + li * 2.0
        rect(s11, lx, leg_y, 0.22, 0.22, col)
        txt(s11, lbl, lx+0.30, leg_y, 1.50, 0.25, 9, color=GRAY)

    # ── SLIDE 12: CERTIFICADOS MINEROS ───────────────────────────────────────
    s12 = prs.slides.add_slide(blank)
    add_bg(s12, bg_contenido)
    slide_header(s12, 'CERTIFICADOS MINEROS')

    txt(s12, 'PRESENTADOS · KPI DE GESTIÓN', 0.45, 1.00, 8.0, 0.28, 10, bold=True, color=TEAL)
    cm_boxes = [(str(cm_tot), 'TOTAL', DARK),
                (str(cm_in),  'IN',    TEAL),
                (str(cm_out), 'OUT',   ROJO if cm_out > 0 else TEAL),
                (fmt_pct(cm_pct), 'KPI', kpi_color(cm_pct))]
    for i, (v, lbl, col) in enumerate(cm_boxes):
        bx = 0.45 + i * 1.72
        rect(s12, bx, 1.35, 1.55, 1.00, WHITE)
        txt(s12, v,   bx, 1.38, 1.55, 0.58, 26, bold=True, color=col,
            align=PP_ALIGN.CENTER, valign='middle')
        txt(s12, lbl, bx, 1.96, 1.55, 0.25,  9, bold=True, color=GRAY,
            align=PP_ALIGN.CENTER)
    txt(s12, 'Límite: 2 días hábiles desde Último Evento',
        0.45, 2.48, 8.0, 0.25, 8.5, color=GRAY)

    txt(s12, 'APROBADOS · TIEMPO DE APROBACIÓN (informativo)',
        0.45, 2.85, 10.0, 0.28, 10, bold=True, color=TEAL)

    rangos_rows = [('0 a 7 días',  r0_7,  pct_r(r0_7),  VERDE),
                   ('8 a 15 días', r8_15, pct_r(r8_15), NARANJA),
                   ('+15 días',    r15,   pct_r(r15),   ROJO)]
    bx2 = 2.40
    bw2 = 9.50
    bh2 = 0.52
    for ri, (lbl, cnt, pct_seg, col) in enumerate(rangos_rows):
        ry2 = 3.27 + ri * 0.85
        txt(s12, lbl, 0.45, ry2, bx2-0.55, bh2, 9, color=DARK, valign='middle')
        rect(s12, bx2, ry2, bw2, bh2, BG_CARD)
        if cnt > 0:
            fw2 = max(0.3, bw2 * pct_seg / 100)
            rect(s12, bx2, ry2, fw2, bh2, col)
            if fw2 > 0.8:
                txt(s12, f"{cnt} exp · {pct_seg}%", bx2+0.10, ry2, fw2-0.1, bh2,
                    9, bold=True, color=WHITE, valign='middle')
            else:
                txt(s12, f"{cnt} exp · {pct_seg}%", bx2+fw2+0.1, ry2, 2.0, bh2,
                    9, bold=True, color=col, valign='middle')
        else:
            txt(s12, 'Sin expedientes', bx2+0.10, ry2, bw2, bh2, 9, color=GRAY, valign='middle')

    txt(s12, f"Total aprobados: {tot_apr} expedientes",
        0.45, 5.60, 6.0, 0.25, 8, color=GRAY)

    # ── SLIDE 13: DESVÍOS CERTIFICADOS MINEROS ──────────────────────────────
    add_desvios_slide(prs, bg_contenido, cm_pre_items,
                      'DESVÍOS · CERTIFICADOS MINEROS', blank, es_cm=True)

    # ── SLIDE 14: CIERRE ─────────────────────────────────────────────────────
    s14 = prs.slides.add_slide(blank)
    cierre_candidates = [
        os.path.join(os.path.dirname(os.path.abspath(__file__)), 'bg_cierre.png'),
        os.path.join(os.path.dirname(os.path.abspath(__file__)), 'bg_cierre.jpg'),
        '/home/claude/bg_cierre.png',
        '/home/claude/bg_cierre.jpg',
    ]
    cierre_bg = next((p for p in cierre_candidates if os.path.exists(p)), None)
    if cierre_bg:
        from io import BytesIO
        with open(cierre_bg, 'rb') as f:
            add_bg(s14, f.read())
    else:
        teal_bg = RGBColor(0x1A, 0x8A, 0x7A)
        teal_dk = RGBColor(0x15, 0x75, 0x67)
        rect(s14, 0, 0, W_IN, H_IN, teal_bg)
        rect(s14, 0, 0, 4.5, H_IN, teal_dk)
        rect(s14, 8.83, 0, 4.5, H_IN, teal_dk)
        rect(s14, 5.88, 1.90, 0.55, 0.55, teal_dk)
        txt(s14, 'A', 5.88, 1.90, 0.55, 0.55, 20, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER, valign='middle')
        txt(s14, 'INTERLOG', 5.30, 2.55, 2.70, 0.40,
            16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s14, 'Comercio Exterior', 5.30, 2.95, 2.70, 0.30,
            10, color=RGBColor(0xC0, 0xE8, 0xE2), align=PP_ALIGN.CENTER)
        for i, line in enumerate([
            'Paseo Colón 505 - 1° Piso - (C1063AEF) - Buenos Aires',
            '0 54 11 5352-3000  –  4362-3055 / 3066  –  4365-5445/5985',
            'info@interlog.com.ar',
        ]):
            txt(s14, line, 2.5, 4.10 + i * 0.28, 8.3, 0.26, 9,
                color=WHITE, align=PP_ALIGN.CENTER)
        txt(s14, 'interlog.com.ar', 2.5, 4.98, 8.3, 0.30,
            11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf, f"KPI_INTERLOG_{mes.replace(' ', '_')}.pptx"


# ── Test ─────────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    from datetime import datetime, timedelta

    def fake_lib(nombre, via, canal, total, out_count):
        return [{'razon': FASA if nombre == 'FASA' else FSM,
                 'nombre': nombre, 'ref': f'R{i}', 'carpeta': f'C{i}',
                 'via': via, 'canal': canal,
                 'f_ofi': datetime(2025, 12, 1), 'f_cancel': datetime(2025, 12, 3),
                 'hs': 2.0, 'limite': 1.0, 'desvio': i < out_count,
                 'desvio_desc': 'Demora Aduana' if i < max(1, out_count//2) else 'Demora Operativa',
                 'parametro': 'ADUANA' if i < max(1, out_count//2) else ('OPERATIVA' if i < out_count else '')}
                for i in range(total)]

    def fake_ofi(nombre, via, total, out_count):
        lim = 2 if nombre == 'FSM' and via == 'MARITIMO' else 1
        return [{'razon': FASA if nombre == 'FASA' else FSM,
                 'nombre': nombre, 'ref': f'O{i}', 'carpeta': f'OC{i}', 'via': via,
                 'f_ofi': datetime(2025, 12, 1), 'f_ult': datetime(2025, 12, 2),
                 'hs': lim+2 if i < out_count else lim-1, 'limite': lim,
                 'desvio': i < out_count, 'desvio_desc': '',
                 'parametro': 'INTERLOG' if i < out_count else ''}
                for i in range(total)]

    lib = (fake_lib('FASA', 'AVION',    'VERDE',   17, 2) +
           fake_lib('FASA', 'AVION',    'NARANJA',  7, 0) +
           fake_lib('FASA', 'AVION',    'ROJO',     2, 0) +
           fake_lib('FASA', 'CAMION',   'VERDE',    5, 1) +
           fake_lib('FASA', 'CAMION',   'NARANJA',  1, 0) +
           fake_lib('FASA', 'CAMION',   'ROJO',     1, 0) +
           fake_lib('FSM',  'AVION',    'VERDE',    9, 2) +
           fake_lib('FSM',  'AVION',    'NARANJA',  6, 4) +
           fake_lib('FSM',  'AVION',    'ROJO',     3, 0) +
           fake_lib('FSM',  'CAMION',   'VERDE',    7, 1) +
           fake_lib('FSM',  'CAMION',   'NARANJA',  1, 1) +
           fake_lib('FSM',  'MARITIMO', 'ROJO',     5, 3))

    ofi = (fake_ofi('FASA', 'AVION',    25, 3) +
           fake_ofi('FSM',  'AVION',    18, 6) +
           fake_ofi('FSM',  'CAMION',    7, 0) +
           fake_ofi('FSM',  'MARITIMO',  5, 0))

    cm_pre = [{'carpeta': f'C{i}', 'exp': f'EXP{i:03}',
               'f_tad': datetime(2025, 12, 4), 'f_ult': datetime(2025, 12, 1),
               'hs': 3 if i < 3 else 1, 'desvio': i < 3,
               'desvio_desc': 'Demora TAD' if i < 3 else '',
               'parametro': 'INTERLOG' if i < 2 else ('ADUANA' if i < 3 else '')}
              for i in range(25)]

    cm_apr = [{'carpeta': f'C{i}', 'exp': f'EXP{i:03}',
               'f_inicio': datetime(2025, 11, 1),
               'f_apro': datetime(2025, 11, 1) + timedelta(days=10),
               'dias': 10, 'rango': '8 a 15' if i < 29 else '+15'}
              for i in range(31)]

    buf, fname = generar_ppt(lib, ofi, cm_pre, cm_apr, mes='DICIEMBRE 2025')
    out = f'/home/claude/{fname}'
    with open(out, 'wb') as f:
        f.write(buf.read())
    print(f'✅ {out}')
